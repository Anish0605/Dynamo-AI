# export.py
import io
from docx import Document
from pptx import Presentation
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib.pagesizes import letter
from fastapi.responses import StreamingResponse

def word(history):
    """Generates a professional Word .docx report."""
    doc = Document()
    doc.add_heading('Dynamo AI Research Report', 0)
    for m in history:
        role = "User" if m['role'] == 'user' else "Dynamo AI"
        doc.add_heading(role, level=1)
        doc.add_paragraph(str(m['content']))
    
    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return StreamingResponse(buf, media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document")

def ppt(history):
    """Generates a PowerPoint .pptx presentation."""
    prs = Presentation()
    for m in history[-5:]:
        if m['role'] == 'assistant':
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "Research Insight"
            slide.placeholders[1].text = str(m['content'])[:700]
            
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return StreamingResponse(buf, media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation")

def pdf(history):
    """Generates a professional PDF report."""
    buf = io.BytesIO()
    pdf_doc = SimpleDocTemplate(buf, pagesize=letter)
    styles = getSampleStyleSheet()
    story = [Paragraph("Dynamo AI Intelligence Report", styles['Title']), Spacer(1, 12)]
    
    for m in history:
        role = "<b>User:</b> " if m['role'] == 'user' else "<b>Dynamo AI:</b> "
        story.append(Paragraph(role + str(m['content']), styles['Normal']))
        story.append(Spacer(1, 12))
        
    pdf_doc.build(story)
    buf.seek(0)
    return StreamingResponse(buf, media_type="application/pdf")
