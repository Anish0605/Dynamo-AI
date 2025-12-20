# main.py - Unified Production (Dynamo AI - #1 AI Research OS Made in India)
from fastapi import FastAPI, HTTPException, UploadFile, File, Form
from fastapi.responses import StreamingResponse, FileResponse
from pydantic import BaseModel
from fastapi.middleware.cors import CORSMiddleware
import os, io, shutil, time, uuid, base64, asyncio, json
import google.generativeai as genai
from groq import Groq
from tavily import TavilyClient
from dotenv import load_dotenv
from pypdf import PdfReader
from PIL import Image
import pandas as pd
import matplotlib.pyplot as plt
import aiohttp

# Document Export Imports
from docx import Document
from pptx import Presentation
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib.pagesizes import letter
import edge_tts

load_dotenv()
app = FastAPI()

# Enable CORS for frontend connectivity
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

# --- API CONFIGURATION ---
GEMINI_KEY = os.getenv("GEMINI_API_KEY")
TAVILY_KEY = os.getenv("TAVILY_API_KEY")
GROQ_KEY = os.getenv("GROQ_API_KEY")

genai.configure(api_key=GEMINI_KEY) if GEMINI_KEY else None
groq_client = Groq(api_key=GROQ_KEY) if GROQ_KEY else None
tavily = TavilyClient(api_key=TAVILY_KEY) if TAVILY_KEY else None

class ChatRequest(BaseModel):
    message: str
    history: list = []
    use_search: bool = True
    deep_dive: bool = False
    fact_check: bool = False
    model: str = "gemini-2.0-flash"
    pdf_context: str = None

class ExportRequest(BaseModel):
    history: list

# --- SYSTEM PROMPTS & IDENTITY ---
SYSTEM_PROMPT = """You are Dynamo AI, the #1 AI Research OS made in India. 
Rules:
1. Always identify as Dynamo AI. Never mention other AI models.
2. Use Markdown for all text responses.
3. For Charts, use Mermaid 'xychart-beta'.
4. For Flowcharts/Mindmaps, use Mermaid 'graph TD' or 'mindmap'.
5. [DEEP DIVE]: If this mode is active, provide 3 detailed sections: 1) Technical/Academic Perspective, 2) Practical/Industry Perspective, and 3) Future/Speculative Trends.
6. [QUIZ]: If asked for a quiz, always return it in a block using ```json_quiz [JSON_DATA] ``` format.
"""

# --- CORE CHAT ENDPOINT ---
@app.post("/chat")
async def chat_endpoint(req: ChatRequest):
    try:
        msg_lower = req.message.lower()
        
        # 1. BACKEND IDENTITY LOGIC (High Priority Interceptor)
        identity_queries = ["who are you", "what is your name", "your name", "who developed you", "who created you"]
        if any(q in msg_lower for q in identity_queries):
            return {
                "type": "text", 
                "content": "My name is **Dynamo AI**, the #1 AI Research OS made in India. I am a highly sophisticated research companion built for data analysis, visual intelligence, and deep research."
            }

        # 2. IMAGE GENERATION LOGIC (System Created + Base64 for Download)
        image_triggers = ["image", "generate", "create", "draw", "visualize"]
        if all(t in msg_lower for t in ["image", "create"]) or all(t in msg_lower for t in ["image", "generate"]):
            prompt = req.message.replace("generate image of", "").replace("create image of", "").strip()
            clean_prompt = prompt.replace(" ", "%20")
            # Using Pollinations as the source but fetching through the backend to handle Base64
            image_url = f"https://image.pollinations.ai/prompt/{clean_prompt}?nologo=true&width=1024&height=1024&seed={uuid.uuid4()}"
            
            async with aiohttp.ClientSession() as session:
                async with session.get(image_url) as resp:
                    if resp.status == 200:
                        image_data = await resp.read()
                        img_b64 = base64.b64encode(image_data).decode('utf-8')
                        return {
                            "type": "image_v2", 
                            "content": f"data:image/jpeg;base64,{img_b64}", 
                            "prompt": prompt
                        }
                    else:
                        return {"type": "text", "content": "Dynamo encountered an issue generating your image. Please try a different prompt."}

        # 3. WEB SEARCH & CONTEXT BUILDING
        context_str = ""
        if req.use_search and not req.pdf_context and tavily:
            try:
                search_res = tavily.search(query=req.message, search_depth="advanced" if req.deep_dive else "basic", max_results=5)
                context_str += f"\n\n[DYNAMO WEB SEARCH RESULTS]:\n{search_res.get('results', [])}\n"
            except: pass
            
        if req.pdf_context:
            context_str += f"\n\n[UPLOADED DOCUMENT CONTEXT]:\n{req.pdf_context[:45000]}\n"

        # 4. MODE HANDLERS
        final_instruction = SYSTEM_PROMPT
        if req.deep_dive:
            final_instruction += "\n[MODE: DEEP DIVE] Perform exhaustive research. Give 3 distinct perspectives: Technical, Industry, and Future."
        if req.fact_check:
            final_instruction += "\n[MODE: FACT CHECK] Verify the claim. Provide sources and a Truth Score (0-100%)."

        # 5. MODEL ROUTING (Gemini 2.0 / LLaMA 3 via Groq)
        if "llama3" in req.model and groq_client:
            msgs = [{"role": "system", "content": final_instruction}]
            if context_str: msgs.append({"role": "system", "content": f"Context: {context_str}"})
            for m in req.history[-8:]:
                msgs.append({"role": "user" if m.get('role') == 'user' else "assistant", "content": m.get('content', '')})
            msgs.append({"role": "user", "content": req.message})
            
            completion = groq_client.chat.completions.create(messages=msgs, model=req.model)
            return {"type": "text", "content": completion.choices[0].message.content}
        else:
            # Default to Gemini 2.0 Flash
            model = genai.GenerativeModel(req.model if "gemini" in req.model else "gemini-2.0-flash")
            prompt_parts = [f"{final_instruction}\nContext Data: {context_str}\n"]
            for m in req.history[-8:]:
                role = "User" if m.get('role') == 'user' else "Dynamo AI"
                prompt_parts.append(f"{role}: {m.get('content')}")
            prompt_parts.append(f"User: {req.message}\nDynamo AI:")
            
            response = model.generate_content("\n".join(prompt_parts))
            return {"type": "text", "content": response.text}

    except Exception as e:
        return {"type": "text", "content": f"⚠️ Dynamo Engine Error: {str(e)}"}

# --- SMART ANALYST (CSV/EXCEL Visualization) ---
@app.post("/analyze-file")
async def analyze_file(file: UploadFile = File(...)):
    try:
        contents = await file.read()
        df = pd.read_csv(io.BytesIO(contents)) if file.filename.endswith('.csv') else pd.read_excel(io.BytesIO(contents))
        
        # Plotting Logic
        plt.figure(figsize=(10, 6))
        df.iloc[:20].select_dtypes(include=['number']).plot(kind='bar', color='#EAB308')
        plt.title(f"Visual Trends: {file.filename}"); plt.tight_layout()
        
        buf = io.BytesIO(); plt.savefig(buf, format='png'); buf.seek(0); plt.close()
        img_b64 = base64.b64encode(buf.getvalue()).decode()
        
        # Dynamic Insight Generation
        summary = df.describe().to_string()
        insight = genai.GenerativeModel("gemini-2.0-flash").generate_content(f"Give 1 single business insight from this summary (max 12 words): {summary}").text
        return {"image": f"data:image/png;base64,{img_b64}", "insight": insight}
    except Exception as e: return {"error": str(e)}

# --- DYNAMO RADIO (2-Person Podcast Scripting) ---
@app.post("/generate-radio")
async def generate_radio(req: ChatRequest):
    prompt = f"Create a 2-person podcast script between Alex and Sam based on this research: {req.message}. Return valid JSON list only: [{{'speaker': 'Alex', 'text': '...'}}, ...]"
    try:
        resp = genai.GenerativeModel("gemini-2.0-flash").generate_content(prompt)
        clean_json = resp.text.strip().replace("```json", "").replace("```", "")
        script = json.loads(clean_json)
        
        audio_segments = []
        for entry in script[:6]:
            voice = "en-IN-PrabhatNeural" if entry["speaker"] == "Alex" else "en-IN-AnanyaNeural"
            temp_path = f"{uuid.uuid4()}.mp3"
            await edge_tts.Communicate(entry["text"], voice).save(temp_path)
            audio_segments.append(temp_path)
            
        # Combine segments
        final_filename = f"DynamoRadio_{int(time.time())}.mp3"
        with open(final_filename, "wb") as f_out:
            for seg in audio_segments:
                with open(seg, "rb") as f_in: f_out.write(f_in.read())
                os.remove(seg)
        return FileResponse(final_filename, media_type="audio/mpeg")
    except Exception as e: return {"error": f"Radio module failed: {str(e)}"}

# --- FULL EXPORT SUITE (Word, PPT, PDF) ---
@app.post("/generate-word")
async def generate_word(req: ExportRequest):
    doc = Document()
    doc.add_heading('Dynamo AI Research Summary', 0)
    for entry in req.history:
        p = doc.add_paragraph()
        p.add_run(f"{'User' if entry['role']=='user' else 'Dynamo AI'}: ").bold = True
        p.add_run(entry['content'])
    buf = io.BytesIO(); doc.save(buf); buf.seek(0)
    return StreamingResponse(buf, media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document")

@app.post("/generate-ppt")
async def generate_ppt(req: ExportRequest):
    prs = Presentation()
    for entry in req.history[-5:]: # Export last 5 research turns
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Dynamo Research Insight"
        slide.placeholders[1].text = entry['content'][:800]
    buf = io.BytesIO(); prs.save(buf); buf.seek(0)
    return StreamingResponse(buf, media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation")

@app.post("/generate-report")
async def generate_pdf(req: ExportRequest):
    buf = io.BytesIO()
    pdf = SimpleDocTemplate(buf, pagesize=letter)
    styles = getSampleStyleSheet()
    story = [Paragraph("Dynamo AI - Research Intelligence Report", styles['Title']), Spacer(1, 12)]
    for entry in req.history:
        story.append(Paragraph(f"<b>{'User' if entry['role']=='user' else 'AI'}:</b>", styles['Normal']))
        story.append(Paragraph(entry['content'], styles['Normal']))
        story.append(Spacer(1, 12))
    pdf.build(story); buf.seek(0)
    return StreamingResponse(buf, media_type="application/pdf")

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=int(os.environ.get("PORT", 5000)))
